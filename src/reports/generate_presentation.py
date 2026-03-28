#!/usr/bin/env python3
"""
src/reports/generate_presentation.py

Generates a 10-slide PowerPoint presentation (dark theme) summarising
the Congressional Stock Analysis pipeline.

Usage:
    python src/reports/generate_presentation.py

Output:
    data/output/summary_presentation.pptx

Prerequisites:
    pip install python-pptx
"""

from __future__ import annotations

import os
import sys
import warnings
from io import BytesIO
from datetime import date
from pathlib import Path

warnings.filterwarnings("ignore")

PROJECT_ROOT = Path(__file__).resolve().parents[2]
os.chdir(PROJECT_ROOT)
sys.path.insert(0, str(PROJECT_ROOT / "src"))

import numpy as np
import pandas as pd

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from model.model_realized import PoliticianTradeModel, Config

# ── constants ─────────────────────────────────────────────────────────────────
OUTPUT_PATH       = "data/output/summary_presentation_v2.pptx"
MODEL_PATH        = "data/output/xgboost_model_realized.json"
THRESHOLD_DEFAULT = 0.60
THRESHOLDS        = [0.50, 0.60, 0.70]

SW, SH = 13.33, 7.5   # slide width / height in inches

# dark theme palette
BG     = "0F1117"
CARD   = "161B22"
TEXT   = "E6EDF3"
MUTED  = "7D8590"
BLUE   = "388BFD"
GREEN  = "3FB950"
ORANGE = "F0883E"
RED    = "F85149"
BORDER = "30363D"

_DARK_RC = {
    "figure.facecolor":  f"#{BG}",
    "axes.facecolor":    f"#{CARD}",
    "axes.edgecolor":    f"#{BORDER}",
    "text.color":        f"#{TEXT}",
    "axes.labelcolor":   f"#{MUTED}",
    "xtick.color":       f"#{MUTED}",
    "ytick.color":       f"#{MUTED}",
    "grid.color":        "#21262D",
    "grid.alpha":        1.0,
    "axes.spines.top":   False,
    "axes.spines.right": False,
    "legend.facecolor":  f"#{CARD}",
    "legend.edgecolor":  f"#{BORDER}",
    "legend.labelcolor": f"#{TEXT}",
}


# ── pptx helpers ──────────────────────────────────────────────────────────────
def _rgb(h: str) -> RGBColor:
    return RGBColor.from_string(h)


def _blank_slide(prs: Presentation):
    layout = prs.slide_layouts[6]   # blank
    slide  = prs.slides.add_slide(layout)
    bg     = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = _rgb(BG)
    return slide


def _slide_title(slide, text: str) -> None:
    """Title text + thin blue rule underneath."""
    tb = slide.shapes.add_textbox(Inches(0.45), Inches(0.18), Inches(12.4), Inches(0.85))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text        = text
    r.font.size   = Pt(32)
    r.font.bold   = True
    r.font.color.rgb = _rgb(TEXT)

    rule = slide.shapes.add_shape(1, Inches(0.45), Inches(1.0), Inches(12.4), Inches(0.04))
    rule.fill.solid()
    rule.fill.fore_color.rgb = _rgb(BLUE)
    rule.line.fill.background()


def _textbox(slide, text, l, t, w, h,
             size=15, color=TEXT, bold=False,
             align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text            = text
    r.font.size       = Pt(size)
    r.font.bold       = bold
    r.font.color.rgb  = _rgb(color)
    return tb


def _bullets(slide, items, l, t, w, h,
             size=15, color=TEXT, bullet="▸",
             header=None, header_color=BLUE):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    if header:
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text           = header
        r.font.size      = Pt(size + 2)
        r.font.bold      = True
        r.font.color.rgb = _rgb(header_color)
        first = False
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(6)
        r = p.add_run()
        r.text           = f"{bullet}  {item}"
        r.font.size      = Pt(size)
        r.font.color.rgb = _rgb(color)
    return tb


def _callout(slide, text, l, t, w, h,
             border=ORANGE, size=17, bold=True, align=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(CARD)
    shape.line.color.rgb      = _rgb(border)
    shape.line.width          = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap        = True
    tf.vertical_anchor  = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text           = text
    r.font.size      = Pt(size)
    r.font.bold      = bold
    r.font.color.rgb = _rgb(border)
    return shape


def _stat_box(slide, label, value, l, t, w=2.8, h=1.4,
              value_color=BLUE):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(CARD)
    shape.line.color.rgb      = _rgb(BORDER)
    shape.line.width          = Pt(1)
    # big number
    tb_val = slide.shapes.add_textbox(Inches(l + 0.15), Inches(t + 0.1),
                                      Inches(w - 0.3), Inches(0.75))
    p = tb_val.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text           = str(value)
    r.font.size      = Pt(26)
    r.font.bold      = True
    r.font.color.rgb = _rgb(value_color)
    # label
    tb_lbl = slide.shapes.add_textbox(Inches(l + 0.15), Inches(t + 0.85),
                                      Inches(w - 0.3), Inches(0.45))
    p2 = tb_lbl.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text           = label
    r2.font.size      = Pt(11)
    r2.font.color.rgb = _rgb(MUTED)


def _pptx_table(slide, headers, rows, l, t, w, h,
                col_widths=None, header_bg=BLUE):
    n_rows = 1 + len(rows)
    n_cols = len(headers)
    tbl = slide.shapes.add_table(n_rows, n_cols,
                                 Inches(l), Inches(t),
                                 Inches(w), Inches(h)).table
    # column widths
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Inches(cw)

    def _cell(r, c, text, bg, fg, size=13, bold=False, align=PP_ALIGN.CENTER):
        cell = tbl.cell(r, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(bg)
        tf = cell.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text           = text
        run.font.size      = Pt(size)
        run.font.bold      = bold
        run.font.color.rgb = _rgb(fg)

    for c, h in enumerate(headers):
        _cell(0, c, h, header_bg, BG, size=13, bold=True)
    for r, row in enumerate(rows):
        row_bg = CARD if r % 2 == 0 else BG
        for c, val in enumerate(row):
            _cell(r + 1, c, str(val), row_bg, TEXT, size=12)


def _img(slide, buf, l, t, w, h=None):
    if h:
        slide.shapes.add_picture(buf, Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        slide.shapes.add_picture(buf, Inches(l), Inches(t), Inches(w))


def _fig_buf(fig, dpi=150):
    buf = BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight", dpi=dpi,
                facecolor=fig.get_facecolor(), edgecolor="none")
    buf.seek(0)
    plt.close(fig)
    return buf


# ── native pptx chart helpers ─────────────────────────────────────────────────
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.oxml.ns import qn as _qn
from lxml import etree as _etree


def _spPr_solid(parent_elm, hex_color):
    """Insert or replace a spPr solid fill on the given XML element."""
    existing = parent_elm.find(_qn('c:spPr'))
    if existing is not None:
        parent_elm.remove(existing)
    sp = _etree.SubElement(parent_elm, _qn('c:spPr'))
    sf = _etree.SubElement(sp, _qn('a:solidFill'))
    sc = _etree.SubElement(sf, _qn('a:srgbClr'))
    sc.set('val', hex_color)
    return sp


def _dark_chart(chart):
    """Apply dark theme to any native pptx chart via XML manipulation."""
    chart_space = chart.element          # c:chartSpace
    chart_elm   = chart_space.find(_qn('c:chart'))
    plot_area   = chart_elm.find(_qn('c:plotArea')) if chart_elm is not None else None

    # chart background (chartSpace spPr)
    _spPr_solid(chart_space, BG)

    # plot area background
    if plot_area is not None:
        _spPr_solid(plot_area, CARD)

    # legend font
    if chart.has_legend:
        try:
            chart.legend.font.color.rgb = _rgb(TEXT)
            chart.legend.font.size = Pt(9)
        except Exception:
            pass

    # axis tick labels
    for attr in ('value_axis', 'category_axis'):
        try:
            ax = getattr(chart, attr)
            ax.tick_labels.font.color.rgb = _rgb(MUTED)
            ax.tick_labels.font.size = Pt(9)
        except Exception:
            pass


def _add_bar(slide, categories, values, l, t, w, h,
             bar_colors=None, title=None, horizontal=True,
             value_label=None, data_labels=False):
    """Single-series bar/column chart with optional per-bar coloring."""
    cd = CategoryChartData()
    cd.categories = list(categories)
    cd.add_series('', [float(v) for v in values])
    chart_type = XL_CHART_TYPE.BAR_CLUSTERED if horizontal else XL_CHART_TYPE.COLUMN_CLUSTERED
    gf = slide.shapes.add_chart(chart_type, Inches(l), Inches(t), Inches(w), Inches(h), cd)
    chart = gf.chart
    chart.has_legend = False
    _dark_chart(chart)
    if bar_colors:
        for pt, color in zip(chart.series[0].points, bar_colors):
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = _rgb(color)
    else:
        chart.series[0].format.fill.solid()
        chart.series[0].format.fill.fore_color.rgb = _rgb(BLUE)
    chart.series[0].format.line.fill.background()
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
        chart.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = _rgb(TEXT)
        chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
    if data_labels:
        chart.plots[0].has_data_labels = True
        chart.plots[0].data_labels.font.color.rgb = _rgb(TEXT)
        chart.plots[0].data_labels.font.size = Pt(8)
    if value_label:
        try:
            chart.value_axis.axis_title.text_frame.text = value_label
            chart.value_axis.axis_title.text_frame.paragraphs[0].runs[0].font.color.rgb = _rgb(MUTED)
            chart.value_axis.has_title = True
        except Exception:
            pass
    return chart


def _add_pie(slide, labels, values, l, t, w, h, colors=None, title=None):
    """Native pie chart."""
    cd = CategoryChartData()
    cd.categories = list(labels)
    cd.add_series('', [float(v) for v in values])
    gf = slide.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(l), Inches(t), Inches(w), Inches(h), cd)
    chart = gf.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.color.rgb = _rgb(TEXT)
    chart.legend.font.size = Pt(9)
    _dark_chart(chart)
    default_colors = [GREEN, RED, ORANGE, BLUE, MUTED]
    colors = colors or default_colors
    for pt, color in zip(chart.series[0].points, colors):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = _rgb(color)
    chart.plots[0].has_data_labels = True
    chart.plots[0].data_labels.font.color.rgb = _rgb(BG)
    chart.plots[0].data_labels.font.size = Pt(12)
    chart.plots[0].data_labels.font.bold = True
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
        chart.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = _rgb(MUTED)
        chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(10)
    return chart


def _add_line(slide, categories, series_dict, l, t, w, h,
              colors=None, title=None, show_legend=True):
    """Multi-series line chart."""
    cd = CategoryChartData()
    cd.categories = list(categories)
    for name, vals in series_dict.items():
        cd.add_series(name, [float(v) if v is not None and not (isinstance(v, float) and np.isnan(v)) else None for v in vals])
    gf = slide.shapes.add_chart(XL_CHART_TYPE.LINE, Inches(l), Inches(t), Inches(w), Inches(h), cd)
    chart = gf.chart
    chart.has_legend = show_legend
    if show_legend:
        chart.legend.position = XL_LEGEND_POSITION.TOP
        chart.legend.font.color.rgb = _rgb(TEXT)
        chart.legend.font.size = Pt(9)
    _dark_chart(chart)
    default_colors = [ORANGE, GREEN, BLUE, MUTED]
    for i, series in enumerate(chart.series):
        color = (colors[i] if colors and i < len(colors) else default_colors[i % len(default_colors)])
        series.format.line.color.rgb = _rgb(color)
        series.format.line.width = Pt(2.0)
        try:
            series.marker.format.fill.solid()
            series.marker.format.fill.fore_color.rgb = _rgb(color)
        except Exception:
            pass
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
        chart.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = _rgb(TEXT)
    return chart


def _add_combo(slide, categories, col_data, line_data, l, t, w, h,
               col_color=BLUE, line_colors=None, show_legend=True):
    """
    Combo chart: clustered columns (primary left axis) + lines (secondary right axis).
    col_data:  {name: [values]}  — exactly one column series
    line_data: {name: [values]}  — one or more line series on secondary axis
    """
    col_name = list(col_data.keys())[0]
    col_vals = list(col_data.values())[0]
    line_colors = line_colors or [ORANGE, GREEN, MUTED]

    cd = CategoryChartData()
    cd.categories = list(categories)
    cd.add_series(col_name, [float(v) for v in col_vals])
    gf = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(l), Inches(t), Inches(w), Inches(h), cd
    )
    chart = gf.chart
    chart_elm = chart.element           # c:chartSpace
    # c:plotArea lives at c:chartSpace > c:chart > c:plotArea
    c_chart_node = chart_elm.find(_qn('c:chart'))
    plot_area = c_chart_node.find(_qn('c:plotArea'))

    # Find the column chart element (barChart with barDir=col or bar3DChart)
    col_chart_elm = None
    for tag in (_qn('c:barChart'), _qn('c:bar3DChart'), _qn('c:columnChart'), _qn('c:lineChart')):
        col_chart_elm = plot_area.find(tag)
        if col_chart_elm is not None:
            break

    # Reassign axis IDs to 101 (cat) / 102 (primary val)
    for old in col_chart_elm.findall(_qn('c:axId')):
        col_chart_elm.remove(old)
    for v in ('101', '102'):
        e = _etree.SubElement(col_chart_elm, _qn('c:axId'))
        e.set('val', v)

    # Build lineChart element pointing to secondary axis 103
    lc = _etree.SubElement(plot_area, _qn('c:lineChart'))
    g = _etree.SubElement(lc, _qn('c:grouping')); g.set('val', 'standard')
    sm = _etree.SubElement(lc, _qn('c:smooth')); sm.set('val', '0')

    markers = ['circle', 'square', 'diamond']
    for si, (sname, svals) in enumerate(line_data.items()):
        ser = _etree.SubElement(lc, _qn('c:ser'))
        ie = _etree.SubElement(ser, _qn('c:idx')); ie.set('val', str(si + 1))
        oe = _etree.SubElement(ser, _qn('c:order')); oe.set('val', str(si + 1))
        tx = _etree.SubElement(ser, _qn('c:tx'))
        v_tx = _etree.SubElement(tx, _qn('c:v')); v_tx.text = sname
        # line style
        sp = _etree.SubElement(ser, _qn('c:spPr'))
        ln = _etree.SubElement(sp, _qn('a:ln')); ln.set('w', '25400')
        sf = _etree.SubElement(ln, _qn('a:solidFill'))
        sc = _etree.SubElement(sf, _qn('a:srgbClr'))
        sc.set('val', line_colors[si] if si < len(line_colors) else ORANGE)
        # marker
        mk = _etree.SubElement(ser, _qn('c:marker'))
        sym = _etree.SubElement(mk, _qn('c:symbol')); sym.set('val', markers[si % len(markers)])
        sz = _etree.SubElement(mk, _qn('c:size')); sz.set('val', '5')
        # values
        val_e = _etree.SubElement(ser, _qn('c:val'))
        nr = _etree.SubElement(val_e, _qn('c:numRef'))
        f_e = _etree.SubElement(nr, _qn('c:f')); f_e.text = ''
        nc = _etree.SubElement(nr, _qn('c:numCache'))
        fc = _etree.SubElement(nc, _qn('c:formatCode')); fc.text = '0.0'
        pc = _etree.SubElement(nc, _qn('c:ptCount')); pc.set('val', str(len(svals)))
        for j, v in enumerate(svals):
            pt = _etree.SubElement(nc, _qn('c:pt')); pt.set('idx', str(j))
            ve = _etree.SubElement(pt, _qn('c:v'))
            ve.text = str(float(v)) if (v is not None and not (isinstance(v, float) and np.isnan(v))) else '0'
    # axis refs for line chart
    for v in ('101', '103'):
        e = _etree.SubElement(lc, _qn('c:axId')); e.set('val', v)

    # Patch existing catAx and valAx IDs
    cat_ax = plot_area.find(_qn('c:catAx'))
    pri_val_ax = plot_area.find(_qn('c:valAx'))
    if cat_ax is not None:
        e = cat_ax.find(_qn('c:axId'))
        if e is not None: e.set('val', '101')
        e2 = cat_ax.find(_qn('c:crossAx'))
        if e2 is not None: e2.set('val', '102')
    if pri_val_ax is not None:
        e = pri_val_ax.find(_qn('c:axId'))
        if e is not None: e.set('val', '102')
        e2 = pri_val_ax.find(_qn('c:crossAx'))
        if e2 is not None: e2.set('val', '101')

    # Add secondary value axis
    def _se(parent, tag, **attrs):
        el = _etree.SubElement(parent, _qn(tag))
        for k, v in attrs.items(): el.set(k, v)
        return el
    sva = _etree.SubElement(plot_area, _qn('c:valAx'))
    _se(sva, 'c:axId', val='103')
    sc_e = _se(sva, 'c:scaling')
    _se(sc_e, 'c:orientation', val='minMax')
    _se(sva, 'c:delete', val='0')
    _se(sva, 'c:axPos', val='r')
    _se(sva, 'c:crossAx', val='101')
    _se(sva, 'c:crosses', val='max')
    _se(sva, 'c:tickLblPos', val='nextTo')

    # Style secondary axis tick labels to be visible in dark theme
    txPr_e = _etree.SubElement(sva, _qn('c:txPr'))
    _etree.SubElement(txPr_e, _qn('a:bodyPr'))
    _etree.SubElement(txPr_e, _qn('a:lstStyle'))
    p_xml = _etree.SubElement(txPr_e, _qn('a:p'))
    pPr_e = _etree.SubElement(p_xml, _qn('a:pPr'))
    defRPr_e = _etree.SubElement(pPr_e, _qn('a:defRPr'))
    defRPr_e.set('sz', '900')
    solidFill_e = _etree.SubElement(defRPr_e, _qn('a:solidFill'))
    srgbClr_e = _etree.SubElement(solidFill_e, _qn('a:srgbClr'))
    srgbClr_e.set('val', TEXT)

    _dark_chart(chart)
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = _rgb(col_color)
    chart.series[0].format.line.fill.background()
    chart.has_legend = show_legend
    if show_legend:
        chart.legend.position = XL_LEGEND_POSITION.TOP
        chart.legend.font.color.rgb = _rgb(TEXT)
        chart.legend.font.size = Pt(9)
    return chart


# ── chart generators (all dark-themed) ───────────────────────────────────────
_parse = PoliticianTradeModel._parse_decimal_comma


def _add_car_hist_native(slide, data, l, t, w, h, title=None, show_median_label=False):
    """CAR distribution as a native pptx column chart (green/red coloring by sign)."""
    bins = np.linspace(-1, 2, 61)
    hist, bin_edges = np.histogram(data.clip(-1, 2), bins=bins)
    bin_centers = [(bin_edges[i] + bin_edges[i+1]) / 2 for i in range(len(hist))]
    cats = [f"{c:.1f}" if i % 15 == 0 else "" for i, c in enumerate(bin_centers)]
    bar_colors = [GREEN if c >= 0 else RED for c in bin_centers]

    cd = CategoryChartData()
    cd.categories = cats
    cd.add_series("Count", [int(v) for v in hist])
    gf = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(l), Inches(t), Inches(w), Inches(h), cd
    )
    chart = gf.chart
    chart.has_legend = False
    _dark_chart(chart)
    for pt, color in zip(chart.series[0].points, bar_colors):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = _rgb(color)
    chart.series[0].format.line.fill.background()
    try:
        chart.plots[0].gap_width = 0
    except Exception:
        pass
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
        chart.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = _rgb(TEXT)
        chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
    if show_median_label:
        med = float(data.median())
        _textbox(slide, f"Median: {med:+.3f}",
                 l=l + w - 2.2, t=t + 0.15, w=2.0, h=0.4,
                 size=11, color=ORANGE, bold=True, align=PP_ALIGN.RIGHT)


def _add_time_trend(slide, proc_df, l, t, w, h):
    """Slide 5: quarterly volume (columns) + % positive (line)."""
    trend = proc_df[["Filed", "realized_car_hybrid"]].copy()
    trend["Filed"] = pd.to_datetime(trend["Filed"], errors="coerce")
    trend["positive"] = (trend["realized_car_hybrid"] > 0).astype(int)
    trend = trend.dropna(subset=["Filed"])
    trend["q"] = trend["Filed"].dt.to_period("Q")
    q = (trend.groupby("q")
         .agg(vol=("positive", "count"), pct=("positive", "mean"))
         .reset_index())
    q = q[q["q"] >= pd.Period("2015Q1")].copy()
    q["qs"] = q["q"].astype(str)
    # Show only Q1 labels, others as empty
    labels = [s if s.endswith("Q1") else "" for s in q["qs"]]
    _add_combo(slide,
               categories=labels,
               col_data={"Trade Volume": q["vol"].tolist()},
               line_data={"% Positive": (q["pct"] * 100).tolist()},
               l=l, t=t, w=w, h=h,
               col_color=BLUE,
               line_colors=[ORANGE],
               show_legend=True)




def _pol_stats(proc_df: pd.DataFrame, min_trades: int = 20) -> pd.DataFrame:
    """Aggregate per-politician stats, filtered to min_trades."""
    id_col = "Name" if "Name" in proc_df.columns else "BioGuideID"
    grp = (proc_df.groupby(id_col)["realized_car_hybrid"]
                  .agg(n="count",
                       pct_pos=lambda s: (s > 0).mean() * 100,
                       med_car=lambda s: s.median() * 100)
                  .reset_index()
                  .rename(columns={id_col: "label"})
                  .query(f"n >= {min_trades}"))
    return grp


def _add_pol_combo(slide, grp, best, n, l, t, w, h, x_min=None, x_max=None):
    """Top or bottom N politicians: combo (hit rate bars + median CAR line)."""
    sub = (grp.sort_values("pct_pos", ascending=not best)
              .head(n)
              .sort_values("pct_pos", ascending=True))
    avg_hit = sub["pct_pos"].mean()
    title   = (f"Top {n} — Best  (group avg {avg_hit:.1f}%)" if best
               else f"Bottom {n} — Worst  (group avg {avg_hit:.1f}%)")
    labels  = [f"{row['label']}  (n={int(row['n'])})" for _, row in sub.iterrows()]

    chart = _add_combo(slide,
                       categories=labels,
                       col_data={"Median CAR %": sub["med_car"].tolist()},
                       line_data={"Hit Rate %":   sub["pct_pos"].tolist()},
                       l=l, t=t, w=w, h=h,
                       col_color=BLUE,
                       line_colors=[ORANGE],
                       show_legend=True)
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    chart.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = _rgb(TEXT)
    chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(11)

    if x_min is not None and x_max is not None:
        try:
            chart.value_axis.minimum_scale = x_min
            chart.value_axis.maximum_scale = x_max
        except Exception:
            pass
    return chart


def _build_flag_chart(slide, proc_df, flags, l, t, w, h):
    """Clustered column: % positive for flag=1 (orange) vs flag=0 (blue)."""
    cats, on_vals, off_vals = [], [], []
    for col, lbl in flags:
        if col not in proc_df.columns:
            continue
        c  = pd.to_numeric(proc_df[col], errors="coerce").fillna(0)
        g1 = proc_df.loc[c == 1, "realized_car_hybrid"]
        g0 = proc_df.loc[c == 0, "realized_car_hybrid"]
        cats.append(lbl)
        on_vals.append(float((g1 > 0).mean() * 100) if len(g1) > 0 else 0.0)
        off_vals.append(float((g0 > 0).mean() * 100) if len(g0) > 0 else 0.0)
    cd = CategoryChartData()
    cd.categories = cats
    cd.add_series("On (flag=1)", on_vals)
    cd.add_series("Off (flag=0)", off_vals)
    gf = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(l), Inches(t), Inches(w), Inches(h), cd
    )
    chart = gf.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.TOP
    chart.legend.font.color.rgb = _rgb(TEXT)
    chart.legend.font.size = Pt(9)
    _dark_chart(chart)
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = _rgb(ORANGE)
    chart.series[0].format.line.fill.background()
    chart.series[1].format.fill.solid()
    chart.series[1].format.fill.fore_color.rgb = _rgb(BLUE)
    chart.series[1].format.line.fill.background()


def _add_committee_lift(slide, proc_df, l, t, w, h):
    """All 9 committee categories."""
    _build_flag_chart(slide, proc_df, [
        ("committee_defense_security",    "Defense &\nSecurity"),
        ("committee_finance_housing",     "Finance &\nHousing"),
        ("committee_fiscal_policy",       "Fiscal\nPolicy"),
        ("committee_energy_environment",  "Energy &\nEnvironment"),
        ("committee_health_labor",        "Health &\nLabor"),
        ("committee_commerce_technology", "Commerce\n& Tech"),
        ("committee_agriculture",         "Agriculture"),
        ("committee_infrastructure",      "Infrastructure"),
        ("committee_oversight",           "Oversight"),
    ], l, t, w, h)


def _add_structural_lift(slide, proc_df, l, t, w, h):
    """Chair, Majority, Lobbied flags."""
    _build_flag_chart(slide, proc_df, [
        ("is_committee_chair",    "Committee\nChair"),
        ("is_committee_majority", "Majority\nParty"),
        ("lobbied_any_90d",       "Lobbied\n(90d)"),
    ], l, t, w, h)


def _add_threshold_sweep(slide, y_test_cont, y_prob, l, t, w, h):
    """Slide 8: N trades (cols) + % positive + median CAR (lines, secondary axis)."""
    thresholds = [t_ / 100 for t_ in range(50, 96, 5)]
    ns, pct_pos_vals, med_car_vals = [], [], []
    for thr in thresholds:
        mask = y_prob >= thr
        n = int(mask.sum())
        ns.append(n)
        if n > 0:
            pct_pos_vals.append(float((y_test_cont[mask] > 0).mean() * 100))
            med_car_vals.append(float(np.median(y_test_cont[mask]) * 100))
        else:
            pct_pos_vals.append(float('nan'))
            med_car_vals.append(float('nan'))
    base_pct = float((y_test_cont > 0).mean() * 100)
    base_med = float(np.median(y_test_cont) * 100)
    lbls = [f"{int(t_*100)}%" for t_ in thresholds]
    _add_combo(slide,
               categories=lbls,
               col_data={"N Trades": ns},
               line_data={
                   f"% Positive (base {base_pct:.1f}%)": pct_pos_vals,
                   f"Median CAR % (base {base_med:.1f}%)": med_car_vals,
               },
               l=l, t=t, w=w, h=h,
               col_color=BLUE,
               line_colors=[ORANGE, GREEN],
               show_legend=True)


def _add_confusion_matrix(slide, y_test, y_prob, l, t, w, h):
    """Confusion matrix rendered as a styled pptx table."""
    from sklearn.metrics import confusion_matrix
    y_pred = (y_prob >= THRESHOLD_DEFAULT).astype(int)
    cm = confusion_matrix(y_test, y_pred)
    tn, fp, fn, tp = cm.ravel()
    total = tn + fp + fn + tp

    # Header row + 2 data rows, 3 cols
    _pptx_table(slide,
                headers=["", "Pred: No beat", "Pred: Beat mkt"],
                rows=[
                    ["Actual: No beat", f"TN = {tn:,}",  f"FP = {fp:,}"],
                    ["Actual: Beat mkt", f"FN = {fn:,}", f"TP = {tp:,}"],
                ],
                l=l, t=t, w=w, h=h,
                col_widths=[w * 0.38, w * 0.31, w * 0.31],
                header_bg=BLUE)


_FEATURE_LABELS = {
    "lobbied_any_90d":                   "Company had active lobbying (90d)",
    "Industry match 3":                  "Tier-3 committee–industry overlap",
    "Industry match 1":                  "Tier-1 committee–industry overlap",
    "Industry match 2":                  "Tier-2 committee–industry overlap",
    "committee_oversight":               "Member of Oversight committee",
    "Chamber":                           "House vs Senate",
    "politician_mean_realized_car_past": "Politician's past avg. return (realized)",
    "committee_energy_environment":      "Member of Energy & Environment committee",
    "all_pol_sells_same_ticker_30d":     "Other politicians selling same stock (30d)",
    "committee_finance_housing":         "Member of Finance & Housing committee",
    "Ticker_Industry":                   "Stock industry sector",
    "politician_mean_car_past":          "Politician's past avg. market-adj. return",
    "committee_health_labor":            "Member of Health & Labor committee",
    "is_committee_majority":             "Politician is in majority party",
    "beta":                              "Stock market sensitivity (beta)",
    "politician_hit_rate_past":          "Politician's past win rate",
    "politician_recent_sells_15d":       "Politician's own recent sells (15d)",
    "stock_momentum_30d":                "Stock price momentum (30d)",
    "stock_momentum_90d":                "Stock price momentum (90d)",
    "stock_volatility_30d":              "Stock volatility (30d)",
    "max_committee_rank":                "Seniority rank on committee",
    "is_committee_chair":                "Politician is committee chair",
    "car_traded_to_filed":               "Return from trade date to disclosure",
}

_FEATURE_CATEGORIES = {
    "politician_hit_rate_past":          "F0883E",
    "politician_mean_car_past":          "F0883E",
    "politician_mean_realized_car_past": "F0883E",
    "politician_recent_sells_15d":       "388BFD",
    "all_pol_sells_same_ticker_30d":     "388BFD",
    "car_traded_to_filed":               "388BFD",
    "beta":                              "4ECDC4",
    "stock_momentum_30d":                "4ECDC4",
    "stock_momentum_90d":                "4ECDC4",
    "stock_volatility_30d":              "4ECDC4",
    "Ticker_Industry":                   "4ECDC4",
    "Chamber":                           "4ECDC4",
    "committee_defense_security":        "A78BFA",
    "committee_finance_housing":         "A78BFA",
    "committee_fiscal_policy":           "A78BFA",
    "committee_energy_environment":      "A78BFA",
    "committee_health_labor":            "A78BFA",
    "committee_commerce_technology":     "A78BFA",
    "committee_agriculture":             "A78BFA",
    "committee_infrastructure":          "A78BFA",
    "committee_oversight":               "A78BFA",
    "is_committee_chair":                "A78BFA",
    "is_committee_majority":             "A78BFA",
    "max_committee_rank":                "A78BFA",
    "Industry match 1":                  "A78BFA",
    "Industry match 2":                  "A78BFA",
    "Industry match 3":                  "A78BFA",
    "lobbied_any_90d":                   "FB923C",
}
_CATEGORY_LEGEND = [
    ("Politician Skill",     "F0883E"),
    ("Trade Signals",        "388BFD"),
    ("Stock Characteristics","4ECDC4"),
    ("Committee / Industry", "A78BFA"),
    ("Lobbying",             "FB923C"),
]


def _add_feature_importance(slide, model, xgb_clf, l, t, w, h):
    """Top 10 feature importances, bars colored by category, with legend."""
    fi  = pd.Series(xgb_clf.feature_importances_, index=model.feature_names).sort_values(ascending=False)
    top = fi.head(10).sort_values(ascending=True)
    labels     = [_FEATURE_LABELS.get(n, n) for n in top.index]
    bar_colors = [_FEATURE_CATEGORIES.get(n, BLUE) for n in top.index]

    _add_bar(slide,
             categories=labels,
             values=top.values.tolist(),
             l=l, t=t, w=w, h=h,
             bar_colors=bar_colors,
             horizontal=True,
             title="Feature Importances (top 10)")

    # Legend
    legend_l = l + w + 0.2
    for i, (cat, color) in enumerate(_CATEGORY_LEGEND):
        _textbox(slide, f"■  {cat}",
                 l=legend_l, t=t + i * 0.6,
                 w=3.1, h=0.5,
                 size=11, color=color, bold=False)




# ── pipeline ──────────────────────────────────────────────────────────────────
def run_pipeline():
    import xgboost as xgb
    cfg = Config()

    print("Loading data …")
    df_full = pd.read_parquet(cfg.data_path)

    print("Calculating sell-pressure features (≈5 min) …")
    df_full = PoliticianTradeModel._calculate_sells_pressure(df_full, "Filed")
    df_full = PoliticianTradeModel._calculate_all_pol_sells_same_ticker(
        df_full, "Filed", window_days=30)

    df = df_full[df_full["Transaction"] == "Purchase"].copy()
    df["Traded"] = pd.to_datetime(df["Traded"], errors="coerce")
    max_traded = df["Traded"].max()
    cutoff = (max_traded - pd.DateOffset(months=12)).strftime("%Y-%m-%d")

    model   = PoliticianTradeModel(cfg, cutoff_date=cutoff, horizon_months=12)
    proc_df = model.preprocess(df)

    train_df = model.time_split(proc_df)[0].reset_index(drop=True)
    test_df  = model.time_split(proc_df)[1].reset_index(drop=True)

    y_test      = test_df[model.target_binary].reset_index(drop=True)
    y_test_cont = test_df[model.target_continuous].reset_index(drop=True)

    X_train = model.prepare_features(train_df, is_training=True).reset_index(drop=True)
    X_test  = model.prepare_features(test_df,  is_training=False).reset_index(drop=True)

    xgb_clf = xgb.XGBClassifier(enable_categorical=True)
    xgb_clf.load_model(MODEL_PATH)
    y_prob = xgb_clf.predict_proba(X_test)[:, 1]

    return model, proc_df, train_df, test_df, y_test, y_test_cont, y_prob, df_full, xgb_clf


# ── slide builders ────────────────────────────────────────────────────────────
def slide_title(prs, proc_df, test_df):
    s = _blank_slide(prs)
    # central title
    tb = s.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(1.4))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Congressional Stock Analysis"
    r.font.size = Pt(46); r.font.bold = True; r.font.color.rgb = _rgb(TEXT)

    tb2 = s.shapes.add_textbox(Inches(1.5), Inches(3.2), Inches(10.3), Inches(0.8))
    p2 = tb2.text_frame.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "Do Congress Members Beat the Market?"
    r2.font.size = Pt(22); r2.font.color.rgb = _rgb(BLUE)

    rule = s.shapes.add_shape(1, Inches(3.5), Inches(4.1), Inches(6.3), Inches(0.04))
    rule.fill.solid(); rule.fill.fore_color.rgb = _rgb(BORDER)
    rule.line.fill.background()

    tb3 = s.shapes.add_textbox(Inches(1.5), Inches(4.3), Inches(10.3), Inches(0.6))
    p3 = tb3.text_frame.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    r3.text = f"Machine Learning on STOCK Act Disclosures  ·  {date.today().strftime('%B %Y')}"
    r3.font.size = Pt(14); r3.font.color.rgb = _rgb(MUTED)



def slide_car_definition(prs):
    s = _blank_slide(prs)
    _slide_title(s, "How We Measure Performance: Cumulative Abnormal Return (CAR)")

    _bullets(s, [
        "Returns are computed as log returns — additive over time, correct for compounding",
        "Market benchmark = CAPM-expected return:  E[R] = Rf + β × (Rm − Rf)",
        "β estimated from 756 trading days of close prices (min 60 days required)",
        "CAR = Actual log return − CAPM-expected log return",
        "Positive CAR → stock earned more than its risk justified (positive alpha)",
        "Negative CAR → stock earned less than CAPM predicted for its level of risk",
        "CAR > 0 does NOT mean the stock beat the index — it means the",
        "  risk-adjusted excess return was positive",
        "Realized CAR hybrid target:",
        "  → actual buy-to-sell log return, if position closed within 12 months",
        "  → 12-month fixed-window CAR from filing date, if held longer or no sell recorded",
        "CAR window starts the day AFTER the disclosure filing date",
    ], l=0.45, t=1.2, w=7.0, h=5.8, size=13)

    _add_bar(s,
             categories=["Stock Log\nReturn", "CAPM Expected\nReturn", "CAR (α)"],
             values=[18.0, 8.0, 10.0],
             l=7.4, t=1.3, w=5.6, h=3.8,
             bar_colors=[BLUE, MUTED, GREEN],
             horizontal=False,
             title="Illustrative example  (hypothetical numbers)")

    _callout(s,
             "Positive CAR = compensated MORE for risk than CAPM predicts\n"
             "Negative CAR = compensated LESS — paid for risk that did not materialise",
             l=7.4, t=5.3, w=5.6, h=1.95, border=ORANGE, size=13)


def slide_distribution(prs, proc_df):
    s = _blank_slide(prs)
    _slide_title(s, "Most Congressional Purchases Underperform the Market")

    data_all = proc_df["realized_car_hybrid"]
    med      = data_all.median()
    pct_pos  = (data_all > 0).mean()
    n_total  = len(data_all)

    # Main histogram (native)
    _add_car_hist_native(s, data_all, l=0.3, t=1.15, w=12.7, h=3.4,
                         title="All trades — Realized CAR distribution  (clipped −1 to +2)",
                         show_median_label=True)

    # Sub-histograms: realized (closed within 12m) vs fixed-window
    if "realized_car" in proc_df.columns:
        realized_mask = proc_df["realized_car"].notna()
    else:
        realized_mask = pd.Series([True] * len(proc_df), index=proc_df.index)
    data_realized = proc_df.loc[realized_mask,  "realized_car_hybrid"]
    data_fixed    = proc_df.loc[~realized_mask, "realized_car_hybrid"]

    _add_car_hist_native(s, data_realized, l=0.3,  t=4.65, w=6.2, h=2.1,
                         title=f"Closed within 12m  (n={len(data_realized):,})",
                         show_median_label=True)
    _add_car_hist_native(s, data_fixed,    l=6.8,  t=4.65, w=6.2, h=2.1,
                         title=f"Held to 12m window  (n={len(data_fixed):,})",
                         show_median_label=True)

    _callout(s,
             f"N = {n_total:,} trades  ·  Median CAR {med:+.3f} ({med*100:+.1f}%)  ·  "
             f"only {pct_pos:.1%} beat the market  ·  Green = positive  ·  Red = negative",
             l=0.45, t=6.85, w=12.4, h=0.5, border=ORANGE, size=12)


def slide_time_trend(prs, proc_df):
    s = _blank_slide(prs)
    _slide_title(s, "Performance Over Time: Macro Events Drive Positive Rates")

    _add_time_trend(s, proc_df, l=0.3, t=1.15, w=12.7, h=4.7)

    _bullets(s, [
        "2020 Q2–Q3: surge in positive rate — COVID recovery bounce lifted most purchases",
        "2022: sharp drop aligned with rate-hike drawdown across the market",
        "Volume spikes suggest increased trading activity around election cycles",
        "Positive rate rarely exceeds 60% — even in the best macro environments",
    ], l=0.45, t=5.9, w=12.3, h=1.45, size=13, color=MUTED)


def slide_eda(prs, proc_df):
    """Slide 6: EDA — three structural comparisons on median CAR."""
    s = _blank_slide(prs)
    _slide_title(s, "EDA: Median CAR Across Key Structural Variables")

    # ── 1. Democrats vs Republicans ──────────────────────────────────────────
    party_col = "Party" if "Party" in proc_df.columns else None
    if party_col:
        party_labels = ["Democrat", "Republican"]
        party_meds   = []
        for p in ["D", "R"]:
            sub = proc_df.loc[proc_df[party_col] == p, "realized_car_hybrid"]
            party_meds.append(float(sub.median() * 100) if len(sub) > 0 else 0.0)
        _add_bar(s,
                 categories=party_labels,
                 values=party_meds,
                 l=0.3, t=1.25, w=3.8, h=4.5,
                 bar_colors=[BLUE if v >= 0 else RED for v in party_meds],
                 horizontal=False,
                 title="Median CAR % — Dem vs Rep",
                 value_label="Median CAR %")

    # ── 2. CAR by disclosure delay ────────────────────────────────────────────
    proc2 = proc_df.copy()
    proc2["Filed_dt"]  = pd.to_datetime(proc2["Filed"],  errors="coerce")
    proc2["Traded_dt"] = pd.to_datetime(proc2["Traded"], errors="coerce")
    proc2["delay"]     = (proc2["Filed_dt"] - proc2["Traded_dt"]).dt.days
    bins_d   = [0, 10, 20, 30, 40, 45, 50, 9999]
    lbls_d   = ["0-10d", "10-20d", "20-30d", "30-40d", "40-45d", "45-50d", "50d+"]
    proc2["bucket"] = pd.cut(proc2["delay"], bins=bins_d, labels=lbls_d, right=False)
    dmed = proc2.groupby("bucket", observed=True)["realized_car_hybrid"].median() * 100
    _add_bar(s,
             categories=list(dmed.index.astype(str)),
             values=[float(v) for v in dmed.values],
             l=4.35, t=1.25, w=4.4, h=4.5,
             bar_colors=[GREEN if v >= 0 else RED for v in dmed.values],
             horizontal=False,
             title="Median CAR % by Disclosure Delay",
             value_label="Median CAR %")

    # ── 3. CAR by industry match tier ─────────────────────────────────────────
    match_cols = [c for c in ["Industry match 1", "Industry match 2", "Industry match 3"]
                  if c in proc_df.columns]
    if match_cols:
        proc3 = proc_df.copy()
        def _tier(row):
            v3 = pd.to_numeric(row.get("Industry match 3", 0), errors="coerce")
            v2 = pd.to_numeric(row.get("Industry match 2", 0), errors="coerce")
            v1 = pd.to_numeric(row.get("Industry match 1", 0), errors="coerce")
            if v3 == 1: return "Tier 3 (Direct)"
            if v2 == 1: return "Tier 2 (Oversight)"
            if v1 == 1: return "Tier 1 (Tangential)"
            return "No Match"
        proc3["tier"] = proc3.apply(_tier, axis=1)
        order = ["No Match", "Tier 1 (Tangential)", "Tier 2 (Oversight)", "Tier 3 (Direct)"]
        tmed  = proc3.groupby("tier")["realized_car_hybrid"].median() * 100
        tmed  = tmed.reindex([x for x in order if x in tmed.index])
        _add_bar(s,
                 categories=list(tmed.index.astype(str)),
                 values=[float(v) for v in tmed.values],
                 l=9.05, t=1.25, w=4.0, h=4.5,
                 bar_colors=[GREEN if v >= 0 else RED for v in tmed.values],
                 horizontal=False,
                 title="Median CAR % by Committee-Industry Match Tier",
                 value_label="Median CAR %")

    _bullets(s, [
        "All metrics: median CAR (consistent throughout the deck)",
        "Tier 3 = Direct jurisdiction  ·  Tier 2 = Oversight  ·  Tier 1 = Tangential",
        "Disclosure delay capped at 45 days by STOCK Act; 50d+ entries reflect late filers",
    ], l=0.3, t=5.9, w=12.5, h=1.45, size=12, color=MUTED)


def slide_top_politicians(prs, proc_df):
    s = _blank_slide(prs)
    _slide_title(s, "Who Consistently Beats — and Loses to — the Market?")

    grp   = _pol_stats(proc_df, min_trades=20)
    top10 = grp.sort_values("pct_pos", ascending=False).head(10)
    bot10 = grp.sort_values("pct_pos", ascending=True).head(10)
    all_car = pd.concat([top10["med_car"], bot10["med_car"]])
    x_min = float(all_car.min()) - 1
    x_max = float(all_car.max()) + 1

    _add_pol_combo(s, grp, best=True,  n=10, l=0.25, t=1.15, w=6.4, h=5.4,
                   x_min=x_min, x_max=x_max)
    _add_pol_combo(s, grp, best=False, n=10, l=6.85, t=1.15, w=6.4, h=5.4,
                   x_min=x_min, x_max=x_max)

    _bullets(s, [
        f"Based on {len(grp)} politicians with ≥ 20 disclosed purchases",
        "Bars = median CAR %  ·  Line = % of trades beating the market (hit rate)",
        "Both charts share the same hit rate axis scale for direct comparison",
    ], l=0.3, t=6.7, w=12.5, h=0.7, size=12, color=MUTED)


def slide_committee_lobbying(prs, proc_df):
    s = _blank_slide(prs)
    _slide_title(s, "Committee Membership Shows Selective Advantage")

    _textbox(s, "Committee categories (all 9)", l=0.3, t=1.15, w=7.9, h=0.35,
             size=11, color=MUTED, bold=True)
    _add_committee_lift(s, proc_df, l=0.3, t=1.5, w=7.9, h=4.6)

    _textbox(s, "Structural role flags", l=8.4, t=1.15, w=4.7, h=0.35,
             size=11, color=MUTED, bold=True)
    _add_structural_lift(s, proc_df, l=8.4, t=1.5, w=4.7, h=4.6)

    _bullets(s, [
        "Orange = flag=1 (on committee / in role / company lobbied)  ·  Blue = flag=0",
        "All 9 committee categories shown — previously only 5 were displayed",
        "Lobbied (90d): company had active lobbying disclosure in 90 days before filing date",
    ], l=0.3, t=6.25, w=12.5, h=1.1, size=12, color=MUTED)


def slide_model_accuracy(prs, test_df, y_test, y_test_cont, y_prob):
    from sklearn.metrics import precision_score, confusion_matrix as sk_cm
    s = _blank_slide(prs)
    _slide_title(s, "Median CAR Turns Positive Only at High Confidence (≥ 85%)")

    _add_threshold_sweep(s, y_test_cont, y_prob, l=0.3, t=1.15, w=9.5, h=4.8)

    # Right: confusion matrix at 0.85
    _textbox(s, "Confusion matrix @ 0.85",
             l=9.9, t=1.2, w=3.2, h=0.45, size=12, color=MUTED, align=PP_ALIGN.CENTER)
    y_pred85 = (y_prob >= 0.85).astype(int)
    cm85     = sk_cm(y_test, y_pred85)
    tn, fp, fn, tp = cm85.ravel()
    _pptx_table(s,
                headers=["", "Pred: No beat", "Pred: Beat"],
                rows=[
                    ["Actual: No beat",  f"TN={tn:,}",  f"FP={fp:,}"],
                    ["Actual: Beat mkt", f"FN={fn:,}",  f"TP={tp:,}"],
                ],
                l=9.9, t=1.65, w=3.2, h=1.8,
                col_widths=[1.22, 1.0, 1.0],
                header_bg=BLUE)

    _textbox(s, f"Test set: {len(y_test):,} trades",
             l=9.9, t=3.55, w=3.2, h=0.35, size=11, color=MUTED, align=PP_ALIGN.CENTER)

    # Threshold comparison table
    _textbox(s, "Threshold comparison",
             l=9.9, t=3.95, w=3.2, h=0.35, size=12, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
    thr_rows = []
    for thr in [0.50, 0.60, 0.85]:
        mask = y_prob >= thr
        n_t  = int(mask.sum())
        if n_t > 0:
            pct  = f"{(y_test_cont[mask] > 0).mean():.1%}"
            med  = f"{np.median(y_test_cont[mask])*100:+.2f}%"
            prec = f"{precision_score(y_test, mask.astype(int), zero_division=0):.1%}"
        else:
            pct = med = prec = "—"
        thr_rows.append([f"{int(thr*100)}%", str(n_t), pct, med, prec])
    _pptx_table(s,
                headers=["Thr", "N", "% Pos", "Med CAR", "Prec"],
                rows=thr_rows,
                l=9.9, t=4.35, w=3.2, h=1.55,
                col_widths=[0.50, 0.50, 0.60, 0.78, 0.60],
                header_bg=BLUE)

    prec85 = precision_score(y_test, y_pred85, zero_division=0)
    n85    = int((y_prob >= 0.85).sum())
    med85  = float(np.median(y_test_cont[y_prob >= 0.85])) * 100
    pct85  = float((y_test_cont[y_prob >= 0.85] > 0).mean()) * 100
    _callout(s,
             f"Threshold 0.85 → {pct85:.0f}% win rate  ·  median CAR {med85:+.2f}%  ·  {n85} trades  ·  {prec85:.1%} precision",
             l=0.3, t=6.1, w=9.5, h=0.85, border=GREEN, size=13)


def slide_feature_importance(prs, model, proc_df, train_df, test_df, y_test, y_prob, xgb_clf):
    from sklearn.metrics import precision_score
    s = _blank_slide(prs)
    _slide_title(s, "The Model: What Drives the Predictions?")

    _add_feature_importance(s, model, xgb_clf, l=0.3, t=1.15, w=8.5, h=5.5)




def slide_conclusions(prs, proc_df, test_df, y_test, y_prob):
    from sklearn.metrics import precision_score
    s = _blank_slide(prs)
    _slide_title(s, "Conclusions & Caveats")

    prec85 = precision_score(y_test, (y_prob >= 0.85).astype(int), zero_division=0)

    _bullets(s, [
        f"Most trades underperform: only {(proc_df['realized_car_hybrid']>0).mean():.1%} of purchases beat the market",
        "Past politician skill and lobbying activity are among the top model features",
        "Committee membership shows selective lift — most pronounced for Chair and Majority Party roles",
        "Tier-3 committee–industry match (direct jurisdiction) shows the strongest committee-related edge",
        f"XGBoost reaches {prec85:.1%} precision at threshold 0.85 — fewer but higher-quality trades",
    ], l=0.45, t=1.2, w=6.5, h=3.6, size=14,
       header="Key Findings", header_color=GREEN)

    _bullets(s, [
        "No transaction costs or market impact modelled",
        "Test set is a single holdout period — results may not generalise across regimes",
        "Skill signals are backward-looking: past performance not guaranteed to continue",
        "CAR is measured from filing date — model already accounts for disclosure lag,",
        "  but real-world execution is constrained to after public disclosure",
        "Dataset covers disclosed trades only — undisclosed activity not captured",
    ], l=7.1, t=1.2, w=6.0, h=3.6, size=14,
       header="Caveats", header_color=RED)

    _bullets(s, [
        "What more lobby data could add:",
        "  1. Lobbying spend amount & trend (currently: binary 90d presence only)",
        "  2. Issue specificity — which bills/topics overlap with politician's committee jurisdiction",
        "  3. Revolving door: lobbyists who previously worked for the same politician/committee",
        "  4. Proximity signal: days between lobbying start date and trade filed date",
    ], l=0.45, t=4.95, w=12.4, h=2.3, size=13,
       header="Potential Data Enrichment (Lobbying Focus)", header_color=BLUE)


def slide_methodology(prs):
    """Methodology: ticker-committee matching, 3-tier industry match, lobbying signal."""
    s = _blank_slide(prs)
    _slide_title(s, "Methodology: How Key Signals Were Constructed")

    _bullets(s, [
        "Committee membership from timestamped YAML snapshots (unitedstates/congress-legislators on GitHub)",
        "For each trade, the snapshot closest in time to the Filed date is used",
        "Ensures membership reflects the politician's actual committees at trade time",
        "Committee codes truncated to 4 characters to normalize across snapshot formats",
        "Result: binary committee-category flags per trade (9 categories total)",
    ], l=0.45, t=1.2, w=12.3, h=2.0, size=13,
       header="1 — Ticker ↔ Committee Matching", header_color=BLUE)

    _bullets(s, [
        "Manual mapping in config/commette_industry_map.yaml: committee → sector → industry",
        "Tier 3 (Strong): committee has DIRECT JURISDICTION — controls budgets, contracts, or legislation",
        "Tier 2 (Medium): committee has OVERSIGHT — sets general rules, taxes, or compliance standards",
        "Tier 1 (Tangential): committee has INDIRECT impact — general legislative influence only",
        "Result: Industry match 1/2/3 binary flags per trade (1 = politician's committee matches stock's industry at that tier)",
    ], l=0.45, t=3.35, w=12.3, h=2.1, size=13,
       header="2 — 3-Tier Committee–Industry Match", header_color=ORANGE)

    _bullets(s, [
        "Source: QuiverQuant lobbying disclosures API, downloaded per traded ticker",
        "lobbied_any_90d = 1 if the company had any active lobbying disclosure in 90 days before trade's Filed date",
        "Currently binary — lobbying spend amount and issue specificity are future enrichment candidates",
    ], l=0.45, t=5.6, w=12.3, h=1.7, size=13,
       header="3 — Lobbying Signal", header_color=GREEN)


def slide_appendix(prs, df_full, proc_df):
    s = _blank_slide(prs)
    _textbox(s, "APPENDIX", l=0.45, t=0.15, w=2.5, h=0.45,
             size=12, color=MUTED, bold=True)
    _slide_title(s, "Dataset Overview")

    purchases = df_full[df_full["Transaction"] == "Purchase"]
    d_min = pd.to_datetime(df_full["Traded"], errors="coerce").min()
    d_max = pd.to_datetime(df_full["Traded"], errors="coerce").max()

    _stat_box(s, "Total Records",    f"{len(df_full):,}",   l=0.45, t=1.3, value_color=BLUE)
    _stat_box(s, "Purchases",        f"{len(purchases):,}", l=3.55, t=1.3, value_color=GREEN)
    _stat_box(s, "Modelled",         f"{len(proc_df):,}",   l=6.65, t=1.3, value_color=ORANGE)
    _stat_box(s, "Politicians",      f"{df_full['BioGuideID'].nunique():,}" if "BioGuideID" in df_full.columns else "—",
              l=9.75, t=1.3, value_color=BLUE)

    _pptx_table(s,
                headers=["Metric", "Value"],
                rows=[
                    ["Date range",                f"{d_min.strftime('%Y-%m-%d')} – {d_max.strftime('%Y-%m-%d')}"],
                    ["Transaction types",         "Purchase / Sale / Sale (Full) / Sale (Partial) / Exchange"],
                    ["Unique tickers (purchases)",f"{purchases['Ticker'].nunique():,}" if "Ticker" in purchases.columns else "—"],
                    ["Parties",                   "D / R / I"],
                    ["Chambers",                  "Representatives / Senate"],
                    ["Target variable",           "realized_car_hybrid > 0 (binary)"],
                    ["CAR windows available",     "1m / 3m / 6m / 9m / 12m + realized"],
                    ["Train / test split",        "80% / 20%  chronological"],
                ],
                l=0.45, t=2.9, w=12.4, h=4.3,
                col_widths=[3.5, 8.9])


# ── main ──────────────────────────────────────────────────────────────────────
def main():
    print("=" * 70)
    print("  CONGRESSIONAL STOCK ANALYSIS — PRESENTATION GENERATOR")
    print("=" * 70)

    (model, proc_df, train_df, test_df,
     y_test, y_test_cont, y_prob,
     df_full, xgb_clf) = run_pipeline()

    print("\nBuilding presentation …")
    prs = Presentation()
    prs.slide_width  = Inches(SW)
    prs.slide_height = Inches(SH)

    print("  Slide 1 — Title …")
    slide_title(prs, proc_df, test_df)

    print("  Slide 2 — CAR definition …")
    slide_car_definition(prs)

    print("  Slide 3 — Distribution …")
    slide_distribution(prs, proc_df)

    print("  Slide 4 — Top/worst politicians …")
    slide_top_politicians(prs, proc_df)

    print("  Slide 5 — Time trend …")
    slide_time_trend(prs, proc_df)

    print("  Slide 6 — EDA …")
    slide_eda(prs, proc_df)

    print("  Slide 7 — Committee & lobbying …")
    slide_committee_lobbying(prs, proc_df)

    print("  Slide 8 — Model accuracy …")
    slide_model_accuracy(prs, test_df, y_test, y_test_cont, y_prob)

    print("  Slide 9 — Feature importance …")
    slide_feature_importance(prs, model, proc_df, train_df, test_df, y_test, y_prob, xgb_clf)

    print("  Slide 10 — Conclusions …")
    slide_conclusions(prs, proc_df, test_df, y_test, y_prob)

    print("  Slide 11 — Methodology …")
    slide_methodology(prs)

    print("  Appendix — Dataset …")
    slide_appendix(prs, df_full, proc_df)

    os.makedirs("data/output", exist_ok=True)
    prs.save(OUTPUT_PATH)
    size_mb = os.path.getsize(OUTPUT_PATH) / 1e6
    print(f"\nPresentation saved → {OUTPUT_PATH}  ({size_mb:.1f} MB)")


if __name__ == "__main__":
    main()
